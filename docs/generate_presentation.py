"""
Generate Vision Capsule Arabic presentation in PPTX and PDF formats.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ── paths ──
DOCS = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(DOCS, "capsule-imgs")
PPTX_OUT = os.path.join(DOCS, "vision-capsule-ar.pptx")
PDF_OUT = os.path.join(DOCS, "vision-capsule-ar.pdf")

# ── colours ──
BG_CHARCOAL = RGBColor(0x1C, 0x19, 0x17)
TEXT_CREAM = RGBColor(0xF5, 0xF0, 0xEB)
ACCENT_BRASS = RGBColor(0xB8, 0x96, 0x5A)
JEGR_ORANGE = RGBColor(0xC8, 0x6B, 0x28)
OVERLAY_DARK = RGBColor(0x0F, 0x0D, 0x0C)

# Hex versions for PDF
BG_HEX = (0x1C/255, 0x19/255, 0x17/255)
TEXT_HEX = (0xF5/255, 0xF0/255, 0xEB/255)
BRASS_HEX = (0xB8/255, 0x96/255, 0x5A/255)
ORANGE_HEX = (0xC8/255, 0x6B/255, 0x28/255)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def img_path(name):
    return os.path.join(IMG, name)


# ─────────────────────────────────────────────
#  Helper: set slide background to solid colour
# ─────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ─────────────────────────────────────────────
#  Helper: add a textbox with RTL Arabic text
# ─────────────────────────────────────────────
def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=TEXT_CREAM, alignment=PP_ALIGN.RIGHT,
                font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    # Set RTL at paragraph level via XML
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    # Set RTL attribute
    pPr = p._p.get_or_add_pPr()
    pPr.set('rtl', '1')
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=18, bold=False, color=TEXT_CREAM,
                          alignment=PP_ALIGN.RIGHT, line_spacing=1.5,
                          font_name="Arial"):
    """Add textbox with multiple lines, each as a separate paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.4)
        pPr = p._p.get_or_add_pPr()
        pPr.set('rtl', '1')
    return txBox


def add_shape_rect(slide, left, top, width, height, fill_color, alpha=None):
    """Add a rectangle shape (for overlays etc). alpha: 0-100 (opacity %)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if alpha is not None:
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        # Navigate: p:sp > p:spPr > a:solidFill > a:srgbClr
        spPr = shape._element.find(f'{{{ns_p}}}spPr')
        if spPr is not None:
            solidFill = spPr.find(f'{{{ns_a}}}solidFill')
            if solidFill is not None:
                srgb = solidFill.find(f'{{{ns_a}}}srgbClr')
                if srgb is not None:
                    alpha_elem = etree.SubElement(srgb, f'{{{ns_a}}}alpha')
                    alpha_elem.set('val', str(int(alpha * 1000)))
    return shape


def add_image_bg(slide, image_file, overlay_alpha=60):
    """Add image as background with dark overlay."""
    pic = slide.shapes.add_picture(
        img_path(image_file), Inches(0), Inches(0), SLIDE_W, SLIDE_H
    )
    # Move picture to back
    sp = pic._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    # Add dark overlay
    overlay = add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H,
                             BG_CHARCOAL, alpha=overlay_alpha)
    return pic, overlay


def add_image_inset(slide, image_file, left, top, width, height=None):
    """Add image as inset with optional height."""
    if height:
        pic = slide.shapes.add_picture(img_path(image_file), left, top, width, height)
    else:
        pic = slide.shapes.add_picture(img_path(image_file), left, top, width=width)
    return pic


def add_accent_line(slide, left, top, width):
    """Add a brass accent line."""
    line = add_shape_rect(slide, left, top, width, Pt(3), ACCENT_BRASS)
    return line


# ═════════════════════════════════════════════
#  BUILD PPTX
# ═════════════════════════════════════════════
def build_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # ── SLIDE 1: Cover ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_CHARCOAL)
    add_image_bg(slide, "p7.jpg", overlay_alpha=55)
    # Title
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                "عيش المستقبل", font_size=60, bold=True, color=ACCENT_BRASS,
                alignment=PP_ALIGN.CENTER)
    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(1),
                "وحدات سكنية كبسولية فاخرة مُصنّعة في كوردستان العراق",
                font_size=24, color=TEXT_CREAM, alignment=PP_ALIGN.CENTER)
    # Logo names
    add_textbox(slide, Inches(2), Inches(4.5), Inches(9), Inches(0.8),
                "VISION CAPSULE  |  JEGR FOR LIGHTING",
                font_size=20, bold=True, color=TEXT_CREAM,
                alignment=PP_ALIGN.CENTER, font_name="Arial")
    # Bottom accent line
    add_accent_line(slide, Inches(4), Inches(5.5), Inches(5))

    # ── SLIDE 2: About ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_CHARCOAL)
    # Image on left side
    add_image_inset(slide, "p10.jpg", Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.5))
    # Title on right
    add_textbox(slide, Inches(6.5), Inches(0.5), Inches(6), Inches(0.8),
                "فيجن كبسول", font_size=42, bold=True, color=ACCENT_BRASS)
    add_accent_line(slide, Inches(9.5), Inches(1.4), Inches(3))
    # Description
    add_textbox(slide, Inches(6.5), Inches(1.8), Inches(6), Inches(1),
                "نُصمّم ونُصنّع وحدات سكنية كبسولية عصرية بالكامل في كوردستان العراق",
                font_size=18, color=TEXT_CREAM)
    # Vision
    add_textbox(slide, Inches(6.5), Inches(3.0), Inches(6), Inches(0.5),
                "الرؤية", font_size=20, bold=True, color=JEGR_ORANGE)
    add_textbox(slide, Inches(6.5), Inches(3.5), Inches(6), Inches(0.8),
                "مستقبل السياحة العصرية — وحدات ذكية تجمع الطبيعة والتكنولوجيا",
                font_size=16, color=TEXT_CREAM)
    # Mission
    add_textbox(slide, Inches(6.5), Inches(4.5), Inches(6), Inches(0.5),
                "الرسالة", font_size=20, bold=True, color=JEGR_ORANGE)
    add_textbox(slide, Inches(6.5), Inches(5.0), Inches(6), Inches(0.8),
                "مساحات إقامة تربط الناس بالطبيعة والثقافة المحلية",
                font_size=16, color=TEXT_CREAM)

    # ── SLIDE 3: 12m Model ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_CHARCOAL)
    add_image_bg(slide, "p3.jpg", overlay_alpha=65)
    # Title
    add_textbox(slide, Inches(6), Inches(0.5), Inches(6.5), Inches(0.8),
                "كبسولة 12 متر", font_size=42, bold=True, color=ACCENT_BRASS)
    add_accent_line(slide, Inches(9.5), Inches(1.4), Inches(3))
    # Features list
    features_12m = [
        "غرفة نوم رئيسية",
        "حمام كامل",
        "مطبخ مجهز",
        "صالة جلوس",
        "نوافذ بانورامية",
        "أرضيات خشبية",
    ]
    feature_text = "\n".join([f"◆  {f}" for f in features_12m])
    add_multiline_textbox(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(4.5),
                          [f"◆  {f}" for f in features_12m],
                          font_size=22, color=TEXT_CREAM)

    # ── SLIDE 4: Floor Plan ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_CHARCOAL)
    add_textbox(slide, Inches(2), Inches(0.3), Inches(9), Inches(0.8),
                "تصميم 12 متر", font_size=36, bold=True, color=ACCENT_BRASS,
                alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(5), Inches(1.1), Inches(3))
    # Large centered image
    add_image_inset(slide, "p6.jpg", Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))

    # ── SLIDE 5: Specs ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_CHARCOAL)
    # Image right side
    add_image_inset(slide, "p5.jpg", Inches(7), Inches(1.5), Inches(5.8), Inches(4.5))
    # Title
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(6), Inches(0.8),
                "جودة بناء عالية", font_size=42, bold=True, color=ACCENT_BRASS)
    add_accent_line(slide, Inches(3), Inches(1.4), Inches(3.5))
    # Specs
    specs = [
        "هيكل معدني صلب",
        "زجاج مزدوج عازل",
        "ألواح ألمنيوم خارجية",
        "عزل حراري متقدم",
        "خشب معالج وتشطيبات",
    ]
    add_multiline_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4),
                          [f"◆  {s}" for s in specs],
                          font_size=22, color=TEXT_CREAM)

    # ── SLIDE 6: Gallery ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_CHARCOAL)
    add_textbox(slide, Inches(2), Inches(0.2), Inches(9), Inches(0.7),
                "فيجن كبسول في الطبيعة", font_size=36, bold=True,
                color=ACCENT_BRASS, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(5), Inches(0.95), Inches(3))
    # 4 images in 2x2 grid
    gallery_imgs = ["p9.jpg", "p8.jpg", "p10.jpg", "p11.jpg"]
    positions = [
        (Inches(0.5), Inches(1.3), Inches(6), Inches(2.8)),
        (Inches(6.8), Inches(1.3), Inches(6), Inches(2.8)),
        (Inches(0.5), Inches(4.3), Inches(6), Inches(2.8)),
        (Inches(6.8), Inches(4.3), Inches(6), Inches(2.8)),
    ]
    for img_name, (l, t, w, h) in zip(gallery_imgs, positions):
        add_image_inset(slide, img_name, l, t, w, h)

    # ── SLIDE 7: Use Cases ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_CHARCOAL)
    # Subtle background
    add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
                "أين تُستخدم فيجن كبسول؟", font_size=42, bold=True,
                color=ACCENT_BRASS, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(4.5), Inches(1.4), Inches(4))
    # 6 use cases in 2 columns x 3 rows
    uses = [
        "المنتجعات السياحية", "المخيمات الفاخرة",
        "المهرجانات", "الاستراحات العائلية",
        "المشاريع الفندقية", "المكاتب المتنقلة",
    ]
    col_x = [Inches(7), Inches(1.5)]  # RTL: right column first
    row_y = [Inches(2.3), Inches(3.6), Inches(4.9)]
    for idx, use in enumerate(uses):
        col = idx % 2
        row = idx // 2
        # Card background
        card = add_shape_rect(slide, col_x[col], row_y[row],
                              Inches(4.8), Inches(1.0),
                              RGBColor(0x2A, 0x25, 0x22))
        card.line.color.rgb = ACCENT_BRASS
        card.line.width = Pt(1)
        add_textbox(slide, col_x[col] + Inches(0.3), row_y[row] + Inches(0.15),
                    Inches(4.2), Inches(0.7),
                    use, font_size=22, bold=True, color=TEXT_CREAM,
                    alignment=PP_ALIGN.CENTER)

    # ── SLIDE 8: Stats ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_CHARCOAL)
    add_textbox(slide, Inches(1), Inches(0.4), Inches(11), Inches(0.8),
                "فيجن كبسول بالأرقام", font_size=42, bold=True,
                color=ACCENT_BRASS, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(4.5), Inches(1.3), Inches(4))
    # Stats in 3x2 grid
    stats = [
        ("12 متر", "الطول"),
        ("100%", "صناعة عراقية"),
        ("جاهزة", "للتسليم"),
        ("مزدوج", "زجاج عازل"),
        ("أسابيع", "مدة التسليم"),
        ("كامل", "تخصيص حسب الطلب"),
    ]
    stat_cols = [Inches(9), Inches(4.8), Inches(0.6)]
    stat_rows = [Inches(2.0), Inches(4.5)]
    for idx, (value, label) in enumerate(stats):
        col = idx % 3
        row = idx // 3
        x = stat_cols[col]
        y = stat_rows[row]
        # Stat card
        card = add_shape_rect(slide, x, y, Inches(3.5), Inches(2.0),
                              RGBColor(0x2A, 0x25, 0x22))
        card.line.color.rgb = ACCENT_BRASS
        card.line.width = Pt(1)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.3),
                    Inches(3.1), Inches(0.8),
                    value, font_size=36, bold=True, color=JEGR_ORANGE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.1),
                    Inches(3.1), Inches(0.6),
                    label, font_size=18, color=TEXT_CREAM,
                    alignment=PP_ALIGN.CENTER)

    # ── SLIDE 9: Smart Features ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_CHARCOAL)
    add_image_bg(slide, "p3.jpg", overlay_alpha=75)
    # English tagline
    add_textbox(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.7),
                "SMART.  SUSTAINABLE.  SILENT.",
                font_size=32, bold=True, color=ACCENT_BRASS,
                alignment=PP_ALIGN.CENTER, font_name="Arial")
    add_accent_line(slide, Inches(4.5), Inches(1.0), Inches(4))
    # Features - left column
    smart_features_r = [
        "نظام تحكم ذكي",
        "طاقة شمسية اختيارية",
        "عزل LYZAN",
        "تكييف GREE",
    ]
    smart_features_l = [
        "زجاج Golden Glass",
        "حمام ARTVIT",
        "إضاءة JEGR",
        "أنظمة تهوية وكهرباء",
    ]
    add_multiline_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(3.5),
                          [f"◆  {f}" for f in smart_features_r],
                          font_size=20, color=TEXT_CREAM)
    add_multiline_textbox(slide, Inches(1), Inches(1.5), Inches(5.5), Inches(3.5),
                          [f"◆  {f}" for f in smart_features_l],
                          font_size=20, color=TEXT_CREAM)
    # Partners bar
    add_shape_rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5),
                   RGBColor(0x2A, 0x25, 0x22), alpha=80)
    add_textbox(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.4),
                "شركاؤنا", font_size=18, bold=True, color=JEGR_ORANGE,
                alignment=PP_ALIGN.CENTER)
    partners = "LYZAN  ·  Golden Glass  ·  GREE  ·  ARTVIT  ·  SIMO  ·  JEGR"
    add_textbox(slide, Inches(1), Inches(6.1), Inches(11), Inches(0.5),
                partners, font_size=18, bold=True, color=TEXT_CREAM,
                alignment=PP_ALIGN.CENTER, font_name="Arial")

    # ── SLIDE 10: How to Order ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_CHARCOAL)
    # Image right
    add_image_inset(slide, "p1.jpg", Inches(7), Inches(1.5), Inches(5.8), Inches(4.5))
    # Title
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(6), Inches(0.8),
                "خطوات الطلب", font_size=42, bold=True, color=ACCENT_BRASS)
    add_accent_line(slide, Inches(3), Inches(1.4), Inches(3.5))
    # 4 steps
    steps = [
        ("١", "تواصل معنا"),
        ("٢", "حدد التخصيصات"),
        ("٣", "التصنيع"),
        ("٤", "التوصيل والتركيب"),
    ]
    for i, (num, text) in enumerate(steps):
        y = Inches(2.2) + Inches(1.2) * i
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(5.3), y, Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = JEGR_ORANGE
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_CREAM
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.name = "Arial"
        # vertical center
        txBody = circle._element.find(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}txBody')
        if txBody is not None:
            bodyPr = txBody.find(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
            if bodyPr is not None:
                bodyPr.set('anchor', 'ctr')
        # Text
        add_textbox(slide, Inches(0.8), y + Inches(0.05), Inches(4.2), Inches(0.6),
                    text, font_size=24, color=TEXT_CREAM)

    # ── SLIDE 11: Contact ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_CHARCOAL)
    add_image_bg(slide, "p11.jpg", overlay_alpha=60)
    # Title
    add_textbox(slide, Inches(1), Inches(1.0), Inches(11), Inches(1),
                "عيش المستقبل اليوم", font_size=52, bold=True,
                color=ACCENT_BRASS, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(4.5), Inches(2.2), Inches(4))
    # Logo names
    add_textbox(slide, Inches(2), Inches(2.8), Inches(9), Inches(0.7),
                "VISION CAPSULE  |  JEGR FOR LIGHTING",
                font_size=22, bold=True, color=TEXT_CREAM,
                alignment=PP_ALIGN.CENTER, font_name="Arial")
    # Contact info
    contact_lines = [
        "كوردستان العراق",
        "واتساب: +964 XXX XXX XXXX",
        "info@visioncapsule.iq",
    ]
    add_multiline_textbox(slide, Inches(3), Inches(4.0), Inches(7), Inches(2.5),
                          contact_lines,
                          font_size=22, color=TEXT_CREAM,
                          alignment=PP_ALIGN.CENTER)

    prs.save(PPTX_OUT)
    print(f"PPTX saved: {PPTX_OUT}")


# ═════════════════════════════════════════════
#  BUILD PDF with ReportLab
# ═════════════════════════════════════════════
def build_pdf():
    from reportlab.lib.pagesizes import landscape
    from reportlab.lib.units import inch, cm
    from reportlab.pdfgen import canvas
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.lib.colors import Color, HexColor
    from PIL import Image as PILImage

    # Register Arabic fonts
    pdfmetrics.registerFont(TTFont('Arial', 'C:/Windows/Fonts/arial.ttf'))
    pdfmetrics.registerFont(TTFont('Arial-Bold', 'C:/Windows/Fonts/arialbd.ttf'))
    pdfmetrics.registerFont(TTFont('Tahoma', 'C:/Windows/Fonts/tahoma.ttf'))
    pdfmetrics.registerFont(TTFont('Tahoma-Bold', 'C:/Windows/Fonts/tahomabd.ttf'))

    page_w, page_h = landscape((7.5 * inch, 13.333 * inch))
    c = canvas.Canvas(PDF_OUT, pagesize=(page_w, page_h))
    c.setTitle("Vision Capsule - Arabic Presentation")

    bg_color = Color(*BG_HEX)
    text_color = Color(*TEXT_HEX)
    brass_color = Color(*BRASS_HEX)
    orange_color = Color(*ORANGE_HEX)
    card_color = Color(0x2A/255, 0x25/255, 0x22/255)

    def draw_bg():
        c.setFillColor(bg_color)
        c.rect(0, 0, page_w, page_h, fill=1, stroke=0)

    def draw_image_bg(image_file, overlay_alpha=0.55):
        """Draw image as background with overlay."""
        ipath = img_path(image_file)
        # Draw image stretched to fill
        c.drawImage(ipath, 0, 0, page_w, page_h,
                    preserveAspectRatio=False, mask='auto')
        # Dark overlay
        c.setFillColor(Color(*BG_HEX, alpha=overlay_alpha))
        c.rect(0, 0, page_w, page_h, fill=1, stroke=0)

    def draw_image_inset(image_file, x, y, w, h):
        ipath = img_path(image_file)
        c.drawImage(ipath, x, y, w, h,
                    preserveAspectRatio=True, mask='auto')

    def draw_accent_line(x, y, w):
        c.setStrokeColor(brass_color)
        c.setLineWidth(2)
        c.line(x, y, x + w, y)

    def draw_text_center(text, y, font_size=18, color=text_color,
                         font='Arial-Bold'):
        c.setFillColor(color)
        c.setFont(font, font_size)
        # Use Arabic reshaping
        shaped = reshape_arabic(text)
        tw = c.stringWidth(shaped, font, font_size)
        c.drawString((page_w - tw) / 2, y, shaped)

    def draw_text_right(text, x_right, y, font_size=18, color=text_color,
                        font='Arial'):
        c.setFillColor(color)
        c.setFont(font, font_size)
        shaped = reshape_arabic(text)
        tw = c.stringWidth(shaped, font, font_size)
        c.drawString(x_right - tw, y, shaped)

    def draw_card(x, y, w, h, border=True):
        c.setFillColor(card_color)
        if border:
            c.setStrokeColor(brass_color)
            c.setLineWidth(1)
            c.roundRect(x, y, w, h, 4, fill=1, stroke=1)
        else:
            c.rect(x, y, w, h, fill=1, stroke=0)

    def reshape_arabic(text):
        """Reshape Arabic text for correct PDF rendering using arabic_reshaper + bidi."""
        try:
            import arabic_reshaper
            from bidi.algorithm import get_display
            reshaped = arabic_reshaper.reshape(text)
            return get_display(reshaped)
        except ImportError:
            # Fallback: reverse the text for basic RTL display
            # Split into Arabic and non-Arabic segments
            return text

    def new_page():
        c.showPage()

    # ── SLIDE 1: Cover ──
    draw_bg()
    draw_image_bg("p7.jpg", 0.55)
    draw_text_center("عيش المستقبل", page_h - 2.2 * inch,
                     font_size=52, color=brass_color, font='Arial-Bold')
    draw_accent_line(page_w/2 - 2*inch, page_h - 2.5*inch, 4*inch)
    draw_text_center("وحدات سكنية كبسولية فاخرة مُصنّعة في كوردستان العراق",
                     page_h - 3.3 * inch, font_size=22, color=text_color,
                     font='Arial')
    draw_text_center("VISION CAPSULE  |  JEGR FOR LIGHTING",
                     page_h - 4.5 * inch, font_size=20, color=text_color,
                     font='Arial-Bold')
    draw_accent_line(page_w/2 - 2*inch, page_h - 5*inch, 4*inch)

    # ── SLIDE 2: About ──
    new_page()
    draw_bg()
    draw_image_inset("p10.jpg", 0.4*inch, 1.5*inch, 5*inch, 4*inch)
    # Right side content
    rx = page_w - 0.8*inch
    draw_text_right("فيجن كبسول", rx, page_h - 1.2*inch,
                    font_size=38, color=brass_color, font='Arial-Bold')
    draw_accent_line(rx - 3*inch, page_h - 1.5*inch, 3*inch)
    draw_text_right("نُصمّم ونُصنّع وحدات سكنية كبسولية عصرية بالكامل في كوردستان العراق",
                    rx, page_h - 2.2*inch, font_size=16, color=text_color)
    draw_text_right("الرؤية", rx, page_h - 3.2*inch,
                    font_size=18, color=orange_color, font='Arial-Bold')
    draw_text_right("مستقبل السياحة العصرية — وحدات ذكية تجمع الطبيعة والتكنولوجيا",
                    rx, page_h - 3.7*inch, font_size=15, color=text_color)
    draw_text_right("الرسالة", rx, page_h - 4.5*inch,
                    font_size=18, color=orange_color, font='Arial-Bold')
    draw_text_right("مساحات إقامة تربط الناس بالطبيعة والثقافة المحلية",
                    rx, page_h - 5.0*inch, font_size=15, color=text_color)

    # ── SLIDE 3: 12m Model ──
    new_page()
    draw_bg()
    draw_image_bg("p3.jpg", 0.65)
    rx = page_w - 0.8*inch
    draw_text_right("كبسولة 12 متر", rx, page_h - 1.2*inch,
                    font_size=38, color=brass_color, font='Arial-Bold')
    draw_accent_line(rx - 3*inch, page_h - 1.5*inch, 3*inch)
    features_12m = [
        "غرفة نوم رئيسية", "حمام كامل", "مطبخ مجهز",
        "صالة جلوس", "نوافذ بانورامية", "أرضيات خشبية",
    ]
    for i, f in enumerate(features_12m):
        draw_text_right(f"◆  {f}", rx, page_h - (2.3 + i*0.55)*inch,
                        font_size=20, color=text_color)

    # ── SLIDE 4: Floor Plan ──
    new_page()
    draw_bg()
    draw_text_center("تصميم 12 متر", page_h - 0.8*inch,
                     font_size=32, color=brass_color, font='Arial-Bold')
    draw_accent_line(page_w/2 - 1.5*inch, page_h - 1.1*inch, 3*inch)
    draw_image_inset("p6.jpg", 0.5*inch, 0.8*inch, page_w - 1*inch, page_h - 2.2*inch)

    # ── SLIDE 5: Specs ──
    new_page()
    draw_bg()
    draw_image_inset("p5.jpg", page_w - 6*inch, 1.5*inch, 5.5*inch, 4*inch)
    rx = 6.5*inch
    draw_text_right("جودة بناء عالية", rx, page_h - 1.2*inch,
                    font_size=38, color=brass_color, font='Arial-Bold')
    draw_accent_line(rx - 3.5*inch, page_h - 1.5*inch, 3.5*inch)
    specs = [
        "هيكل معدني صلب", "زجاج مزدوج عازل",
        "ألواح ألمنيوم خارجية", "عزل حراري متقدم",
        "خشب معالج وتشطيبات",
    ]
    for i, s in enumerate(specs):
        draw_text_right(f"◆  {s}", rx, page_h - (2.3 + i*0.55)*inch,
                        font_size=20, color=text_color)

    # ── SLIDE 6: Gallery ──
    new_page()
    draw_bg()
    draw_text_center("فيجن كبسول في الطبيعة", page_h - 0.6*inch,
                     font_size=32, color=brass_color, font='Arial-Bold')
    draw_accent_line(page_w/2 - 1.5*inch, page_h - 0.9*inch, 3*inch)
    gallery = ["p9.jpg", "p8.jpg", "p10.jpg", "p11.jpg"]
    gw = (page_w - 1.5*inch) / 2
    gh = (page_h - 2.5*inch) / 2
    positions = [
        (0.5*inch, page_h - 1.3*inch - gh),
        (0.5*inch + gw + 0.5*inch, page_h - 1.3*inch - gh),
        (0.5*inch, page_h - 1.3*inch - 2*gh - 0.3*inch),
        (0.5*inch + gw + 0.5*inch, page_h - 1.3*inch - 2*gh - 0.3*inch),
    ]
    for img_name, (gx, gy) in zip(gallery, positions):
        draw_image_inset(img_name, gx, gy, gw - 0.2*inch, gh - 0.15*inch)

    # ── SLIDE 7: Use Cases ──
    new_page()
    draw_bg()
    draw_text_center("أين تُستخدم فيجن كبسول؟", page_h - 1.0*inch,
                     font_size=38, color=brass_color, font='Arial-Bold')
    draw_accent_line(page_w/2 - 2*inch, page_h - 1.3*inch, 4*inch)
    uses = [
        "المنتجعات السياحية", "المخيمات الفاخرة",
        "المهرجانات", "الاستراحات العائلية",
        "المشاريع الفندقية", "المكاتب المتنقلة",
    ]
    card_w = 4.5*inch
    card_h = 0.9*inch
    col_x = [page_w - 1*inch - card_w, 1*inch]
    row_y = [page_h - 2.5*inch, page_h - 3.8*inch, page_h - 5.1*inch]
    for idx, use in enumerate(uses):
        col = idx % 2
        row = idx // 2
        cx = col_x[col]
        cy = row_y[row]
        draw_card(cx, cy, card_w, card_h)
        draw_text_center(use, cy + 0.3*inch, font_size=20, color=text_color,
                         font='Arial-Bold')

    # ── SLIDE 8: Stats ──
    new_page()
    draw_bg()
    draw_text_center("فيجن كبسول بالأرقام", page_h - 1.0*inch,
                     font_size=38, color=brass_color, font='Arial-Bold')
    draw_accent_line(page_w/2 - 2*inch, page_h - 1.3*inch, 4*inch)
    stats = [
        ("12 متر", "الطول"),
        ("100%", "صناعة عراقية"),
        ("جاهزة", "للتسليم"),
        ("مزدوج", "زجاج عازل"),
        ("أسابيع", "مدة التسليم"),
        ("كامل", "تخصيص حسب الطلب"),
    ]
    stat_w = 3.8*inch
    stat_h = 2*inch
    stat_cols = [page_w - 0.8*inch - stat_w,
                 page_w/2 - stat_w/2,
                 0.8*inch]
    stat_rows = [page_h - 1.8*inch - stat_h, page_h - 1.8*inch - 2*stat_h - 0.3*inch]
    for idx, (value, label) in enumerate(stats):
        col = idx % 3
        row = idx // 3
        sx = stat_cols[col]
        sy = stat_rows[row]
        draw_card(sx, sy, stat_w, stat_h)
        # Value
        c.setFillColor(orange_color)
        c.setFont('Arial-Bold', 32)
        shaped_val = reshape_arabic(value)
        vw = c.stringWidth(shaped_val, 'Arial-Bold', 32)
        c.drawString(sx + (stat_w - vw)/2, sy + stat_h - 0.8*inch, shaped_val)
        # Label
        c.setFillColor(text_color)
        c.setFont('Arial', 16)
        shaped_lab = reshape_arabic(label)
        lw = c.stringWidth(shaped_lab, 'Arial', 16)
        c.drawString(sx + (stat_w - lw)/2, sy + 0.4*inch, shaped_lab)

    # ── SLIDE 9: Smart Features ──
    new_page()
    draw_bg()
    draw_image_bg("p3.jpg", 0.75)
    draw_text_center("SMART.  SUSTAINABLE.  SILENT.", page_h - 0.8*inch,
                     font_size=30, color=brass_color, font='Arial-Bold')
    draw_accent_line(page_w/2 - 2*inch, page_h - 1.1*inch, 4*inch)
    smart_r = [
        "نظام تحكم ذكي", "طاقة شمسية اختيارية",
        "عزل LYZAN", "تكييف GREE",
    ]
    smart_l = [
        "زجاج Golden Glass", "حمام ARTVIT",
        "إضاءة JEGR", "أنظمة تهوية وكهرباء",
    ]
    rx = page_w - 0.8*inch
    for i, f in enumerate(smart_r):
        draw_text_right(f"◆  {f}", rx, page_h - (1.7 + i*0.5)*inch,
                        font_size=18, color=text_color)
    lx = 6*inch
    for i, f in enumerate(smart_l):
        draw_text_right(f"◆  {f}", lx, page_h - (1.7 + i*0.5)*inch,
                        font_size=18, color=text_color)
    # Partners bar
    c.setFillColor(Color(0x2A/255, 0x25/255, 0x22/255, alpha=0.85))
    c.rect(0.5*inch, 0.5*inch, page_w - 1*inch, 1.4*inch, fill=1, stroke=0)
    draw_text_center("شركاؤنا", 1.5*inch, font_size=16,
                     color=orange_color, font='Arial-Bold')
    draw_text_center("LYZAN  ·  Golden Glass  ·  GREE  ·  ARTVIT  ·  SIMO  ·  JEGR",
                     0.9*inch, font_size=17, color=text_color, font='Arial-Bold')

    # ── SLIDE 10: How to Order ──
    new_page()
    draw_bg()
    draw_image_inset("p1.jpg", page_w - 6*inch, 1.2*inch, 5.5*inch, 4.5*inch)
    rx = 6.5*inch
    draw_text_right("خطوات الطلب", rx, page_h - 1.2*inch,
                    font_size=38, color=brass_color, font='Arial-Bold')
    draw_accent_line(rx - 3.5*inch, page_h - 1.5*inch, 3.5*inch)
    steps = [
        ("١", "تواصل معنا"),
        ("٢", "حدد التخصيصات"),
        ("٣", "التصنيع"),
        ("٤", "التوصيل والتركيب"),
    ]
    for i, (num, text) in enumerate(steps):
        y = page_h - (2.3 + i*1.1)*inch
        # Orange circle
        c.setFillColor(orange_color)
        c.circle(rx - 0.1*inch, y + 0.15*inch, 0.3*inch, fill=1, stroke=0)
        c.setFillColor(text_color)
        c.setFont('Arial-Bold', 20)
        shaped_num = reshape_arabic(num)
        nw = c.stringWidth(shaped_num, 'Arial-Bold', 20)
        c.drawString(rx - 0.1*inch - nw/2, y + 0.03*inch, shaped_num)
        # Text
        draw_text_right(text, rx - 0.8*inch, y + 0.05*inch,
                        font_size=22, color=text_color)

    # ── SLIDE 11: Contact ──
    new_page()
    draw_bg()
    draw_image_bg("p11.jpg", 0.60)
    draw_text_center("عيش المستقبل اليوم", page_h - 1.8*inch,
                     font_size=46, color=brass_color, font='Arial-Bold')
    draw_accent_line(page_w/2 - 2*inch, page_h - 2.2*inch, 4*inch)
    draw_text_center("VISION CAPSULE  |  JEGR FOR LIGHTING",
                     page_h - 3.0*inch, font_size=20, color=text_color,
                     font='Arial-Bold')
    contact = [
        "كوردستان العراق",
        "واتساب: +964 XXX XXX XXXX",
        "info@visioncapsule.iq",
    ]
    for i, line in enumerate(contact):
        draw_text_center(line, page_h - (4.0 + i*0.55)*inch,
                         font_size=20, color=text_color, font='Arial')

    c.save()
    print(f"PDF saved: {PDF_OUT}")


if __name__ == "__main__":
    build_pptx()
    build_pdf()
